#!/usr/bin/env python3
"""
OpenClaw Introduction Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(26, 26, 46)
CYAN = RGBColor(0, 217, 255)
RED = RGBColor(255, 107, 107)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    # Content
    y = Inches(1.5)
    for item in content_items:
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.333), Inches(0.8))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        y += Inches(0.9)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    y = Inches(2.2)
    for item in left_items:
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(5.5), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "  • " + item
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        y += Inches(0.7)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    y = Inches(2.2)
    for item in right_items:
        box = slide.shapes.add_textbox(Inches(7), y, Inches(5.5), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "  • " + item
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        y += Inches(0.7)
    
    return slide

# Create slides
# Slide 1: Title
add_title_slide(prs, "OpenClaw 🦞", "自托管 AI 代理网关平台")

# Slide 2: What is OpenClaw
add_content_slide(prs, "什么是 OpenClaw？", [
    "自托管网关，连接聊天应用与 AI 代理",
    "支持 WhatsApp、Telegram、Discord、iMessage 等",
    "在本地运行，数据完全掌控",
    "MIT 开源许可，社区驱动",
    "口号：'EXFOLIATE! EXFOLIATE!' — 来自一只太空龙虾 🦞"
])

# Slide 3: Core Features
add_two_column_slide(prs, "核心功能", 
    "主要特性", [
        "多通道网关（单一进程支持多平台）",
        "插件扩展（Mattermost 等）",
        "多代理路由",
        "媒体支持（图片、音频、文档）"
    ],
    "技术能力", [
        "Web 控制界面",
        "iOS/Android 节点配对",
        "Canvas 渲染",
        "相机/语音工作流"
    ]
)

# Slide 4: Architecture
add_content_slide(prs, "系统架构", [
    "聊天应用 + 插件 → Gateway（网关）",
    "Gateway → Pi 代理 / CLI / Web 控制界面 / 移动节点",
    "网关是会话、路由和通道连接的单一事实来源",
    "每个发送者独立会话，保证隐私隔离"
])

# Slide 5: Who Should Use It
add_two_column_slide(prs, "目标用户",
    "适合人群", [
        "开发者",
        "高级用户",
        "需要个人 AI 助手的人",
        "注重数据隐私的用户"
    ],
    "使用场景", [
        "随时随地消息访问",
        "不依赖托管服务",
        "完全控制数据",
        "自定义 AI 行为"
    ]
)

# Slide 6: Quick Start
add_content_slide(prs, "快速开始", [
    "1. 安装：npm install -g openclaw@latest",
    "2. 引导：openclaw onboard --install-daemon",
    "3. 配对 WhatsApp：openclaw channels login",
    "4. 启动网关：openclaw gateway --port 18789",
    "5. 访问控制 UI：http://127.0.0.1:18789/"
])

# Slide 7: Configuration
add_content_slide(prs, "配置选项", [
    "配置文件：~/.openclaw/openclaw.json",
    "默认使用 bundled Pi 二进制（RPC 模式）",
    "可按发送者设置白名单（allowFrom）",
    "群组消息可设置提及规则（requireMention）",
    "支持自定义 API 提供商和模型"
])

# Slide 8: Community & Resources
add_content_slide(prs, "社区与资源", [
    "官方文档：https://docs.openclaw.ai",
    "GitHub: https://github.com/openclaw/openclaw",
    "Discord 社区：https://discord.com/invite/clawd",
    "技能市场：https://clawhub.com",
    "MIT 开源许可"
])

# Slide 9: Summary
add_title_slide(prs, "OpenClaw 🦞", "让 AI 助手无处不在，同时保持数据掌控权")

# Save presentation
output_path = "/Users/mike/.openclaw/workspace/OpenClaw_Introduction.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
